import os
from pptx import Presentation
from pptx.util import Inches, Cm
import os
# 画像フォルダのパス
image_folders = [
    'C:\\Users\\user\\python\\ppt\\sample_images\\images1',
    'C:\\Users\\user\\python\\ppt\\sample_images\\images2',
    'C:\\Users\\user\\python\\ppt\\sample_images\\images3'
]
print(image_folders)
template_file = 'test_template.pptx'
output_file = 'output.pptx'
# パワーポイントプレゼンテーションを作成
presentation = Presentation(template_file)
# 画像フォルダ内の画像ファイルを取得
for image_folder in image_folders:
    image_files = [f for f in os.listdir(image_folder) if f.endswith(('.jpg', '.jpeg', '.png', '.gif'))]
    print(image_files)
    # スライド上に画像を配置する位置とサイズの設定
    left = Cm(1)  # 左端の位置（センチメートル単位）
    top = Cm(2)  # 上端の位置（センチメートル単位）
    width = Cm(5)  # 画像の幅（センチメートル単位）
    height = Cm(4)  # 画像の高さ（センチメートル単位）
    space = Cm(1)  # 画像間のスペース（センチメートル単位）
    # 画像をスライドに配置する
    for i, image_file in enumerate(image_files):
        print('i ', i)
        print(image_file)
        if i % 9 == 0:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # レイアウト6を使用するスライドを作成
            left = Cm(1)  # 左端の位置をリセット
            top = Cm(2)  # 上端の位置をリセット
        image_path = os.path.join(image_folder, image_file)
        image = slide.shapes.add_picture(image_path, left, top, width, height)
        left += width + space  # 次の画像の左端位置を更新
        k=(i+1) % 3
        if k == 0:
            left=Cm(1)
            top += height + space  # 次の行の上端位置を更新
# パワーポイントファイルを保存
presentation.save(output_file)
# 保存したパワーポイントファイルを開く（Windows環境の場合）
os.startfile(output_file)